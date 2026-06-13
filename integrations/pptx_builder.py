"""Генерация PowerPoint (.pptx) для презентаций."""

from __future__ import annotations

import io
import os
import re
from typing import Optional


def _safe_title(text: str, max_len: int = 80) -> str:
    t = re.sub(r"[\\/*?:\[\]]", "", (text or "").strip())
    return (t[:max_len] + "…") if len(t) > max_len else t or "Presentation"


def _parse_slides(task_text: str, response: str) -> list[dict]:
    slides = []
    lines = (response or "").split("\n")
    current = {"title": "Введение", "bullets": []}
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if re.match(r"^#{1,3}\s+", line) or re.match(r"^\d+[\.)]\s+", line):
            if current["bullets"] or current["title"] != "Введение":
                slides.append(current)
            title = re.sub(r"^#{1,3}\s+", "", line)
            title = re.sub(r"^\d+[\.)]\s+", "", title)
            current = {"title": title[:100], "bullets": []}
        elif line.startswith(("- ", "• ", "* ")):
            current["bullets"].append(line.lstrip("-•* ").strip()[:200])
        elif len(current["bullets"]) < 6:
            current["bullets"].append(line[:200])
    if current["bullets"] or current["title"]:
        slides.append(current)
    if not slides:
        slides = [
            {"title": _safe_title(task_text), "bullets": [task_text[:300]]},
            {"title": "Ключевые пункты", "bullets": [(response or task_text)[:400][:200]]},
            {"title": "Следующие шаги", "bullets": ["Обсудить с командой", "Уточнить детали", "Запустить реализацию"]},
        ]
    return slides[:12]


def build_pptx_bytes(task_text: str, response: str = "", title: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    deck_title = _safe_title(title or task_text)
    slides_data = _parse_slides(task_text, response)

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = deck_title
    try:
        slide.placeholders[1].text = "AI Team Room · Ника"
    except Exception:
        pass

    for sd in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sd.get("title", "Slide")[:100]
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = sd.get("bullets") or ["—"]
        for i, b in enumerate(bullets[:8]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def save_pptx(task_text: str, response: str, out_dir: str, filename: str = "presentation.pptx") -> str:
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, filename)
    with open(path, "wb") as f:
        f.write(build_pptx_bytes(task_text, response))
    return path
